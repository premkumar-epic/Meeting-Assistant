from __future__ import annotations

import os
from pathlib import Path
from docx import Document
from pptx import Presentation

CONTEXT_DIR = Path(__file__).resolve().parents[2] / "context"
OUTPUT_FILE = CONTEXT_DIR / "context_dump.txt"


def extract_docx(file_path: Path) -> str:
    out = []
    out.append(f"\n=========================================\nFILE: {file_path.name}\n=========================================\n")
    try:
        doc = Document(str(file_path))
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                out.append(text)
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells]
                row_text = [t for t in row_text if t]
                if row_text:
                    out.append(" | ".join(row_text))
    except Exception as e:
        out.append(f"[ERROR PARSING DOCX: {e}]")
    return "\n".join(out)


def extract_pptx(file_path: Path) -> str:
    out = []
    out.append(f"\n=========================================\nFILE: {file_path.name}\n=========================================\n")
    try:
        prs = Presentation(str(file_path))
        for slide_idx, slide in enumerate(prs.slides):
            out.append(f"--- Slide {slide_idx + 1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    out.append(shape.text.strip())
    except Exception as e:
        out.append(f"[ERROR PARSING PPTX: {e}]")
    return "\n".join(out)


def main() -> None:
    all_text = []
    for file in sorted(CONTEXT_DIR.iterdir()):
        if file.suffix == ".docx":
            print(f"Extracting {file.name}...")
            all_text.append(extract_docx(file))
        elif file.suffix == ".pptx":
            print(f"Extracting {file.name}...")
            all_text.append(extract_pptx(file))
            
    with OUTPUT_FILE.open("w", encoding="utf-8") as f:
        f.write("\n\n".join(all_text))
        
    print(f"Extraction complete! Content saved to {OUTPUT_FILE}")


if __name__ == "__main__":
    main()
